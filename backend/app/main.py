from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from dotenv import load_dotenv
from io import BytesIO
import re
import os
import json

from groq import Groq
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from .database import Base, engine, SessionLocal
from . import models, schemas
from .auth import hash_password
from app.core.rag_engine import RAGEngine
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from typing import List, Dict, Any
from collections import defaultdict
from reportlab.platypus import SimpleDocTemplate, Spacer, Preformatted


from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4

from app.routers import gap_analysis
from app.routers import transformation



load_dotenv()

# ------------------ APP & DB ------------------
Base.metadata.create_all(bind=engine)
app = FastAPI(title="TeachAssist Backend")

app.include_router(gap_analysis.router)
app.include_router(transformation.router)

# ------------------ CORS ------------------
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:4200"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ------------------ DB DEP ------------------
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# ------------------ LLM + RAG ------------------
rag = RAGEngine()
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))
GROQ_MODEL = "llama-3.1-8b-instant"

# =========================================================
# FILE TEXT EXTRACTION
# =========================================================
def extract_text_from_file(upload_file: UploadFile) -> str:
    raw = upload_file.file.read()
    filename = upload_file.filename.lower()

    if filename.endswith(".pdf"):
        reader = PdfReader(BytesIO(raw))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if filename.endswith(".pptx"):
        prs = Presentation(BytesIO(raw))
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slides_text.append(shape.text)
        return "\n".join(slides_text)

    if filename.endswith(".docx"):
        doc = Document(BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if filename.endswith(".txt"):
        return raw.decode("utf-8", errors="ignore")

    raise HTTPException(
        status_code=400,
        detail="Unsupported file format. Upload PDF, PPTX, DOCX, or TXT."
    )

# ------------------ CLEAN OUTPUT ------------------
def clean_output(text: str) -> str:
    text = re.sub(r"\*{1,3}", "", text)
    text = text.replace("•", "-").replace("###", "")
    return text.strip()

# ------------------ TEXT CHUNKING ------------------
def chunk_text(text: str, size=800, overlap=100):
    chunks, start = [], 0
    while start < len(text):
        chunks.append(text[start:start + size])
        start += size - overlap
    return chunks

# ------------------ EXTRACT MARKS FROM PROMPT ------------------
def extract_marks_config(prompt: str) -> Dict[str, int]:
    """Extract marks configuration from the prompt"""
    marks = {
        'mcq': 1,
        'short': 5,
        'long': 10,
        'scenario': 10,
        'task': 4  # Default for assignment tasks
    }
    
    # Try to extract marks from prompt
    mcq_match = re.search(r'MCQ.*?(\d+)\s*marks', prompt, re.IGNORECASE)
    short_match = re.search(r'Short.*?(\d+)\s*marks', prompt, re.IGNORECASE)
    long_match = re.search(r'Long.*?(\d+)\s*marks', prompt, re.IGNORECASE)
    scenario_match = re.search(r'Scenario.*?(\d+)\s*marks', prompt, re.IGNORECASE)
    task_match = re.search(r'MARKS_PER_TASK:\s*(\d+)', prompt, re.IGNORECASE)
    
    if mcq_match:
        marks['mcq'] = int(mcq_match.group(1))
    if short_match:
        marks['short'] = int(short_match.group(1))
    if long_match:
        marks['long'] = int(long_match.group(1))
    if scenario_match:
        marks['scenario'] = int(scenario_match.group(1))
    if task_match:
        marks['task'] = int(task_match.group(1))
    
    return marks

# ------------------ PARSE EXAM STRUCTURE ------------------
def parse_exam_for_export(content: str, marks_config: Dict[str, int], exam_type: str) -> Dict[str, Any]:
    """Parse exam content for structured export"""
    lines = content.split('\n')
    sections = []
    current_section = None
    current_qtype = None
    answer_key_started = False
    
    questions_part = []
    answers_part = []
    tasks = []  # For assignments
    
    for line in lines:
        if '=== ANSWER KEY ===' in line:
            answer_key_started = True
            continue
        
        if answer_key_started:
            answers_part.append(line)
        else:
            questions_part.append(line)
    
    # Handle ASSIGNMENT type
    if exam_type == 'assignment':
        task_num = 0
        for line in questions_part:
            line_stripped = line.strip()
            if line_stripped.startswith('Task '):
                task_num += 1
                tasks.append(line_stripped)
        
        return {
            'type': 'assignment',
            'tasks': tasks,
            'answers': '\n'.join(answers_part),
            'marks_config': marks_config
        }
    
    # Handle QUIZ and MID/FINAL types
    for line in questions_part:
        line_stripped = line.strip()
        
        if not line_stripped:
            continue
        
        # Detect section headers
        if line_stripped.startswith('SECTION'):
            if current_section:
                sections.append(current_section)
            
            section_name = line_stripped.split(':')[0].strip()
            current_section = {
                'name': section_name,
                'mcqs': [],
                'shorts': [],
                'longs': [],
                'scenarios': [],
                'total_marks': 0
            }
            current_qtype = None
        
        # Detect question type headers
        elif 'Multiple Choice' in line_stripped:
            current_qtype = 'mcqs'
        elif 'Short Answer' in line_stripped:
            current_qtype = 'shorts'
        elif 'Long Answer' in line_stripped:
            current_qtype = 'longs'
        elif 'Scenario' in line_stripped and current_section:
            current_qtype = 'scenarios'
        elif line_stripped and current_section and current_qtype:
            # Add question to current type
            if line_stripped[0].isdigit() or line_stripped.startswith('Task'):
                current_section[current_qtype].append(line_stripped)
    
    if current_section:
        sections.append(current_section)
    
    # Calculate total marks for each section
    for section in sections:
        section['total_marks'] = (
            len(section['mcqs']) * marks_config.get('mcq', 1) +
            len(section['shorts']) * marks_config.get('short', 5) +
            len(section['longs']) * marks_config.get('long', 10) +
            len(section['scenarios']) * marks_config.get('scenario', 10)
        )
    
    return {
        'type': 'exam',
        'sections': sections,
        'answers': '\n'.join(answers_part),
        'marks_config': marks_config
    }

# ------------------ FORMATTED DOCX EXPORT ------------------
# ------------------ FORMATTED DOCX EXPORT ------------------
def generate_docx(content: str, include_answers: bool, prompt: str = "", exam_type: str = "quiz"):
    """
    VERY SIMPLE EXPORT:
    - Ignore exam_type / marks / structure
    - Just write the raw exam.content line by line into a DOCX
    - This guarantees that NOTHING is dropped
    """
    doc = Document()

    # Optional: simple heading
    # doc.add_heading('TeachAssist Export', level=1)

    for line in content.splitlines():
        # Empty line => blank paragraph (keeps spacing)
        if line.strip() == "":
            doc.add_paragraph("")
        else:
            doc.add_paragraph(line)

    return doc

def generate_pdf(content: str, include_answers: bool, prompt: str = "", exam_type: str = "quiz"):
    """
    SUPER SIMPLE PDF EXPORT

    - No parsing / sections
    - Dumps the ENTIRE `content` string to PDF
    - Preserves line breaks, wraps long lines to page width
    """
    buffer = BytesIO()

    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36,
    )

    styles = getSampleStyleSheet()
    body_style = ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=10,
        leading=13,
        wordWrap="LTR",   # wrap long lines
    )

    # Escape minimal HTML and convert newlines to <br/> for Paragraph
    safe_text = (
        content.replace("&", "&amp;")
               .replace("<", "&lt;")
               .replace(">", "&gt;")
    )
    safe_text = safe_text.replace("\r\n", "\n").replace("\r", "\n")
    safe_text = safe_text.replace("\n", "<br/>")

    story = [Paragraph(safe_text, body_style)]
    doc.build(story)

    buffer.seek(0)
    return buffer



# ------------------ SIMPLE RAW PDF EXPORT ------------------
@app.get("/api/download/{exam_id}")
def download_exam(
    exam_id: int,
    format: str,
    include_answers: bool = True,
    db: Session = Depends(get_db)
):
    exam = db.query(models.GeneratedExam).get(exam_id)
    if not exam:
        raise HTTPException(404, "Exam not found")

    exam_type = exam.exam_type if hasattr(exam, "exam_type") else "quiz"
    prompt = ""

    if format == "docx":
        doc = generate_docx(exam.content, include_answers, prompt, exam_type)
        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": "attachment; filename=exam.docx"}
        )

    if format == "pdf":
        pdf = generate_pdf(exam.content, include_answers, prompt, exam_type)
        return StreamingResponse(
            pdf,
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=exam.pdf"}
        )

    raise HTTPException(400, "Invalid format")


# ------------------ LLaMA GENERATION ------------------
def generate_with_llama(prompt: str) -> str:
    """Generate with focused system prompt"""
    response = groq_client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[
            {
                "role": "system",
                "content": "You are a precise university exam generator. Follow instructions exactly. Generate questions ONLY from provided context. Create realistic workplace scenarios when requested. Count questions carefully before outputting."
            },
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.2,
        max_tokens=8196
    )
    return response.choices[0].message.content
# ===================== ASSIGNMENT HELPERS (add above the endpoint) =====================
ASSIGN_TOTAL_RE = re.compile(r"TOTAL_TASKS:\s*(\d+)", re.IGNORECASE)
ASSIGN_SCEN_RE  = re.compile(r"SCENARIO_TASKS:\s*(\d+)", re.IGNORECASE)

def extract_assignment_requirements(prompt: str) -> tuple[int, int]:
    t = ASSIGN_TOTAL_RE.search(prompt)
    s = ASSIGN_SCEN_RE.search(prompt)
    total = int(t.group(1)) if t else 0
    scenario = int(s.group(1)) if s else 0
    return total, scenario


def _parse_assignment_tasks(questions: str):
    """
    Split the questions-part into:
      - prefix_lines: everything before Task 1
      - tasks: list of {index, header, body}
    We accept headers like: 'Task 1:', 'Task 1 (5 marks):', etc.
    """
    lines = questions.splitlines()
    prefix_lines: list[str] = []
    tasks: list[dict] = []

    current_header: str | None = None
    current_index: int | None = None
    current_body_lines: list[str] = []

    for line in lines:
        m = re.match(r"\s*(Task\s+(\d+)[^:]*:)", line)
        if m:
            # starting a new Task block
            if current_header is None:
                # first task -> everything before is prefix
                prefix_lines = []
            else:
                # store previous task
                tasks.append({
                    "index": current_index,
                    "header": current_header,
                    "body": "\n".join(current_body_lines).rstrip()
                })

            current_header = m.group(1).rstrip()   # only up to the colon
            current_index = int(m.group(2))

            # anything after the colon on the same line becomes the first body line
            remainder = line[m.end():].rstrip()
            current_body_lines = [remainder] if remainder else []
        else:
            if current_header is None:
                prefix_lines.append(line)
            else:
                current_body_lines.append(line)

    # flush last task
    if current_header is not None:
        tasks.append({
            "index": current_index,
            "header": current_header,
            "body": "\n".join(current_body_lines).rstrip()
        })

    return prefix_lines, tasks


def _normalize_to_scenario(body: str) -> str:
    """
    Convert the body of a task into:
        Scenario: ...
        Task: ...
    If it's already in proper Scenario+Task format, keep as is.
    """
    has_scenario = re.search(r"(?im)^\s*Scenario\s*:", body) is not None
    has_task     = re.search(r"(?im)^\s*Task\s*:", body)     is not None

    # already proper scenario → just tidy
    if has_scenario and has_task:
        return body.strip() + "\n"

    lines = [ln for ln in body.strip().splitlines() if ln.strip()]
    first_line = lines[0] if lines else ""
    rest_text  = "\n".join(lines[1:]).strip()

    scenario_stub = (
        "Scenario: Consider a realistic workplace/project situation "
        "relevant to the uploaded lecture content."
    )
    task_line = (
        f"Task: {first_line.strip()}" if first_line
        else "Task: Respond to the scenario."
    )

    if rest_text:
        task_line += f"\n{rest_text}"

    return f"{scenario_stub}\n{task_line}\n"


def force_assignment_scenarios(text: str, total_tasks: int, scenario_tasks: int) -> str:
    """
    Ensure the last `scenario_tasks` tasks are SCENARIO style.
    - If SCENARIO_TASKS == TOTAL_TASKS → ALL tasks become scenarios.
    - Only modifies the questions part (before '=== ANSWER KEY ===').
    """
    if total_tasks <= 0 or scenario_tasks <= 0:
        return text

    parts = text.split("=== ANSWER KEY ===", 1)
    questions = parts[0]
    answers   = parts[1] if len(parts) > 1 else ""

    prefix_lines, tasks = _parse_assignment_tasks(questions)
    if not tasks:
        # nothing recognized as tasks, bail out
        return text

    present_indices = [t["index"] for t in tasks]

    # Clamp scenario_tasks to actually existing tasks
    scenario_tasks = min(scenario_tasks, len(present_indices))

    # We always convert the LAST N tasks into scenarios
    target_indices = present_indices[-scenario_tasks:]
    target_set = set(target_indices)

    # Rebuild the questions part
    out_lines: list[str] = []
    out_lines.extend(prefix_lines)

    for t in tasks:
        body = t["body"]
        if t["index"] in target_set:
            body = _normalize_to_scenario(body)

        block = t["header"] + "\n"
        if body:
            block += body
        out_lines.append(block)

    normalized_questions = "\n".join(out_lines).strip()

    if answers:
        return f"{normalized_questions}\n\n=== ANSWER KEY ===\n{answers.lstrip()}"
    return normalized_questions

def trim_assignment_tasks(text: str, total_tasks: int) -> str:
    """
    Keep only Task 1..TOTAL_TASKS in both QUESTIONS and ANSWER KEY.
    Drop any Task N where N > TOTAL_TASKS.
    """
    if not total_tasks or total_tasks <= 0:
        return text

    parts = text.split("=== ANSWER KEY ===", 1)
    questions = parts[0]
    answers   = parts[1] if len(parts) > 1 else ""

    # ---- trim questions part using _parse_assignment_tasks ----
    prefix_lines, tasks = _parse_assignment_tasks(questions)
    if not tasks or len(tasks) <= total_tasks:
        # nothing to trim
        return text

    kept_tasks = tasks[:total_tasks]

    out_q_lines: list[str] = []
    out_q_lines.extend(prefix_lines)

    for t in kept_tasks:
        block = t["header"]
        body  = t["body"]
        out_q_lines.append(block)
        if body:
            out_q_lines.append(body)
        # blank line between tasks (optional but keeps layout nice)
        out_q_lines.append("")

    normalized_questions = "\n".join(out_q_lines).strip()

    # ---- trim answer key (if present) ----
    if not answers:
        return normalized_questions

    filtered_answer_lines: list[str] = []
    for line in answers.splitlines():
        m = re.match(r"\s*Task\s+(\d+)\b", line)
        if m:
            idx = int(m.group(1))
            if idx > total_tasks:
                # skip answers for extra tasks
                continue
        filtered_answer_lines.append(line)

    normalized_answers = "\n".join(filtered_answer_lines).strip()

    return f"{normalized_questions}\n\n=== ANSWER KEY ===\n{normalized_answers}\n"

# =========================================================
# GENERATE EXAM (MAIN ENDPOINT)
# =========================================================
@app.post("/api/generate-exam", response_model=schemas.ExamOut)
async def generate_exam(
    exam_type: str = Form(...),
    prompt: str = Form(...),
    teacher_prompt: str = Form(""),
    files: List[UploadFile] = File(...),
    db: Session = Depends(get_db)
):
    try:
        # normalize exam type once
        exam_type_lower = exam_type.lower()

        # ------------------ Ensure teacher ------------------
        teacher = db.query(models.Teacher).first()
        if not teacher:
            teacher = models.Teacher(
                username="default_teacher",
                password_hash=hash_password("1234")
            )
            db.add(teacher)
            db.commit()
            db.refresh(teacher)

        # ------------------ Extract text from files ------------------
        all_documents = []
        for idx, file in enumerate(files):
            text = extract_text_from_file(file)
            if text.strip():
                all_documents.append({
                    "doc_id": idx + 1,
                    "content": text
                })

        if not all_documents:
            raise HTTPException(400, "Empty lecture content")

        # ------------------ Chunk + Store in RAG ------------------
        all_chunks = []
        all_metadata = []

        for doc in all_documents:
            chunks = chunk_text(doc["content"])
            all_chunks.extend(chunks)
            all_metadata.extend(
                [{"teacher_id": teacher.id, "doc_id": doc["doc_id"]}] * len(chunks)
            )

        rag.add_documents(texts=all_chunks, metadatas=all_metadata)

        # ------------------ Retrieve balanced context from RAG ------------------
        doc_chunks_map = defaultdict(list)
        docs_count = len(all_documents)
        TOTAL_K = 12
        k_per_doc = max(1, TOTAL_K // docs_count)

        for doc in all_documents:
            chunks = rag.search(
                query="Generate exam questions strictly from this document",
                top_k=k_per_doc,
                filter={"doc_id": doc["doc_id"]}
            )
            doc_chunks_map[doc["doc_id"]].extend(chunks)

        if not doc_chunks_map:
            raise HTTPException(400, "No relevant content retrieved from RAG")

        # Build context with document separation
        context_parts = []
        for doc_id in sorted(doc_chunks_map.keys()):
            chunks = doc_chunks_map[doc_id]
            joined_chunks = "\n\n".join(chunks)
            context_parts.append(
                f"""===== DOCUMENT {doc_id} =====
{joined_chunks}"""
            )

        context = "\n\n".join(context_parts)

   # ------------------ Build final prompt (FRONTEND DRIVEN) ------------------

        final_prompt = f"""
You are an AI exam generator for teachers.

You MUST strictly follow the instructions provided in the UI PROMPT.
You MUST generate questions ONLY from the lecture context below.

================ LECTURE CONTEXT =================
{context}

================ UI PROMPT (STRICT RULES) =================
{prompt}

================ TEACHER NOTES =================
{teacher_prompt}
"""

        # Single LLaMA call for ALL exam types
        raw_exam_text = generate_with_llama(final_prompt)
        exam_text = clean_output(raw_exam_text)

                # Extra post-processing ONLY for assignments
        if exam_type_lower == "assignment":
            total_tasks, scenario_tasks = extract_assignment_requirements(prompt)

            # 1) Enforce the number of tasks (drop Task 6,7,8,9, …)
            if total_tasks > 0:
                exam_text = trim_assignment_tasks(exam_text, total_tasks)

            # 2) Force the last N tasks to be scenarios (if any)
            if scenario_tasks > 0:
                exam_text = force_assignment_scenarios(
                    exam_text,
                    total_tasks,
                    scenario_tasks
                )

        # ------------------ Save to database ------------------
        exam = models.GeneratedExam(
            teacher_id=teacher.id,
            exam_type=exam_type,   # original string ('quiz', 'assignment', 'midterm')
            content=exam_text
        )
        db.add(exam)
        db.commit()
        db.refresh(exam)

        return exam

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ------------------ DOWNLOAD ENDPOINT ------------------
@app.get("/api/download/{exam_id}")
def download_exam(
    exam_id: int,
    format: str,
    include_answers: bool = True,
    db: Session = Depends(get_db)
):
    exam = db.query(models.GeneratedExam).get(exam_id)
    if not exam:
        raise HTTPException(404, "Exam not found")

    exam_type = exam.exam_type if hasattr(exam, "exam_type") else "quiz"
    prompt = ""

    if format == "docx":
        doc = generate_docx(exam.content, include_answers, prompt, exam_type)
        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": "attachment; filename=exam.docx"}
        )

    if format == "pdf":
        pdf = generate_pdf(exam.content, include_answers, prompt, exam_type)
        return StreamingResponse(
            pdf,
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=exam.pdf"}
        )

    raise HTTPException(400, "Invalid format")
